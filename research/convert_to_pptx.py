#!/usr/bin/env python3
"""
Convert presentation markdown to PowerPoint PPTX
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def parse_slides(markdown_content):
    """Parse markdown into slide sections"""
    # Split by slide markers (---\n\n## Slide)
    slide_pattern = r'## Slide \d+:'
    slides = re.split(slide_pattern, markdown_content)

    # Get slide titles
    titles = re.findall(slide_pattern, markdown_content)

    parsed_slides = []
    for i, slide_content in enumerate(slides[1:]):  # Skip intro text
        slide_data = {
            'title': titles[i].replace('## Slide ', '').replace(':', '').strip() if i < len(titles) else '',
            'content': '',
            'visual': '',
            'text': '',
            'speaker_notes': ''
        }

        # Extract sections from each slide
        visual_match = re.search(r'\*\*VISUAL:\*\*\s*(.*?)(?=\*\*TEXT:|SPEAKER NOTES:|$)', slide_content, re.DOTALL)
        text_match = re.search(r'\*\*TEXT:\*\*\s*(.*?)(?=\*\*SPEAKER NOTES:|$)', slide_content, re.DOTALL)
        notes_match = re.search(r'\*\*SPEAKER NOTES:\*\*\s*(.*?)(?=\n---|\Z)', slide_content, re.DOTALL)

        if visual_match:
            slide_data['visual'] = visual_match.group(1).strip()
        if text_match:
            slide_data['text'] = text_match.group(1).strip()
        if notes_match:
            slide_data['speaker_notes'] = notes_match.group(1).strip()

        # If no TEXT section, use visual as content
        if not slide_data['text'] and slide_data['visual']:
            slide_data['content'] = slide_data['visual']
        else:
            slide_data['content'] = slide_data['text']

        parsed_slides.append(slide_data)

    return parsed_slides

def clean_markdown(text):
    """Remove markdown formatting for PowerPoint"""
    # Remove bold/italic
    text = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'__(.+?)__', r'\1', text)
    text = re.sub(r'_(.+?)_', r'\1', text)

    # Remove links but keep text
    text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)

    # Clean up code blocks
    text = re.sub(r'```[\w]*\n', '', text)
    text = re.sub(r'```', '', text)

    # Remove emojis from text (optional - keep them for now)

    return text.strip()

def extract_bullet_points(text):
    """Extract bullet points from text"""
    lines = text.split('\n')
    bullets = []
    for line in lines:
        line = line.strip()
        if line.startswith('- ') or line.startswith('* '):
            bullets.append(clean_markdown(line[2:]))
        elif line.startswith('• '):
            bullets.append(clean_markdown(line[2:]))
        elif line and not line.startswith('#'):
            bullets.append(clean_markdown(line))
    return [b for b in bullets if b]  # Filter empty

def create_presentation(slides_data, output_file):
    """Create PowerPoint presentation from parsed slides"""
    prs = Presentation()

    # Set slide dimensions (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for idx, slide_data in enumerate(slides_data):
        # Determine slide layout based on content
        if idx == 0:  # Title slide
            slide_layout = prs.slide_layouts[0]  # Title Slide
        elif 'table' in slide_data['content'].lower() or '|' in slide_data['content']:
            slide_layout = prs.slide_layouts[5]  # Blank for tables
        else:
            slide_layout = prs.slide_layouts[1]  # Title and Content

        slide = prs.slides.add_slide(slide_layout)

        # Add title
        if slide.shapes.title:
            title_text = slide_data['title']
            if not title_text and idx == 0:
                title_text = "Bounded Iterative Vibing"
            slide.shapes.title.text = title_text

            # Format title
            title_frame = slide.shapes.title.text_frame
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Deep blue

        # Add content
        content_text = slide_data['content']
        if content_text:
            # Try to find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    content_placeholder = shape
                    break

            if content_placeholder:
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                text_frame.word_wrap = True

                # Extract bullets and add to slide
                bullets = extract_bullet_points(content_text)
                if bullets:
                    for bullet in bullets[:8]:  # Limit to 8 bullets per slide
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.font.size = Pt(18)
                        p.space_after = Pt(12)
                else:
                    # Add as regular text
                    p = text_frame.paragraphs[0]
                    p.text = clean_markdown(content_text[:500])  # Limit length
                    p.font.size = Pt(20)
            else:
                # Add text box manually
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(4.5)

                txBox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = txBox.text_frame
                text_frame.word_wrap = True

                bullets = extract_bullet_points(content_text)
                if bullets:
                    for bullet in bullets[:8]:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.font.size = Pt(18)
                        p.space_after = Pt(12)
                else:
                    p = text_frame.paragraphs[0]
                    p.text = clean_markdown(content_text[:500])
                    p.font.size = Pt(20)

        # Add speaker notes
        if slide_data['speaker_notes']:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = clean_markdown(slide_data['speaker_notes'])

    # Save presentation
    prs.save(output_file)
    print(f"✅ Presentation saved to: {output_file}")
    print(f"📊 Total slides: {len(slides_data)}")

def main():
    # Read markdown file
    input_file = '/home/user/dotai/research/presentation-bounded-iterative-vibing.md'
    output_file = '/home/user/dotai/research/bounded-iterative-vibing-presentation.pptx'

    print("📖 Reading presentation markdown...")
    with open(input_file, 'r', encoding='utf-8') as f:
        markdown_content = f.read()

    print("🔍 Parsing slides...")
    slides_data = parse_slides(markdown_content)

    if not slides_data:
        print("❌ No slides found! Check markdown format.")
        return

    print(f"✅ Found {len(slides_data)} slides")
    print("🎨 Creating PowerPoint presentation...")

    create_presentation(slides_data, output_file)

    print("\n✨ Done! Your presentation is ready.")
    print(f"📍 Location: {output_file}")
    print("\n💡 Next steps:")
    print("   1. Open in PowerPoint/Keynote/Google Slides")
    print("   2. Adjust layouts and add visuals")
    print("   3. Apply your brand theme")
    print("   4. Add diagrams and charts")

if __name__ == '__main__':
    main()
